from typing import List, Callable, Dict

from pptx.chart.data import CategoryChartData
from pptx.compat import to_unicode
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.shapes.graphfrm import GraphicFrame
from pptx.text.text import TextFrame

from .type import KeyValueDataType, CategoryDataType, KVKeys, ChartDataType, ChartDataTypes, Text
from .utils import find_shape, find_shape_by_slide_layout


def replace_paragraphs_text(paragraph, new_text):
    # copy from https://github.com/scanny/python-pptx/issues/285#issuecomment-300273686
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    if paragraph.runs:
        for idx, run in enumerate(paragraph.runs):
            if idx == 0:
                continue
            p.remove(run._r)
    else:
        paragraph.add_run()
    paragraph.runs[0].text = new_text


def replace_run_style(style_run, target_run):
    target_run.font.name = style_run.font.name
    target_run.font.size = style_run.font.size
    target_run.font.bold = style_run.font.bold
    target_run.font.italic = style_run.font.italic
    target_run.font.underline = style_run.font.underline
    target_run.font.language_id = style_run.font.language_id

    if color_type := style_run.font.color.type:
        if color_type == MSO_COLOR_TYPE.RGB:
            target_run.font.color.rgb = style_run.font.color.rgb
        elif color_type == MSO_COLOR_TYPE.SCHEME:
            target_run.font.color.theme_color = style_run.font.color.theme_color
            target_run.font.color.brightness = style_run.font.color.brightness


def replace_paragraph_style(style_paragraph, target_paragraph):
    style_paragraph.alignment = target_paragraph.alignment
    style_paragraph.level = target_paragraph.level
    style_paragraph.line_spacing = target_paragraph.line_spacing
    style_paragraph.space_after = target_paragraph.space_after
    style_paragraph.space_before = target_paragraph.space_before


def clear_text_frame(text_frame: TextFrame):
    # remove all paragraphs
    for p in text_frame._txBody.p_lst[1:]:
        text_frame._txBody.remove(p)


def replace_text_frame(text_frame: TextFrame, new_text: Text):
    clear_text_frame(text_frame)
    style_paragraph = text_frame.paragraphs[0]
    style_run = style_paragraph.runs[0] if len(style_paragraph.runs) else style_paragraph.add_run()

    split_text = to_unicode(f"{new_text}").split('\n')
    if split_text:
        replace_paragraphs_text(text_frame.paragraphs[0], split_text[0])
        for text in split_text[1:]:
            p = text_frame.add_paragraph()
            replace_paragraph_style(style_paragraph, p)
            replace_paragraphs_text(p, text)
            replace_run_style(style_run, p.runs[0])


def replace_table_data_key_value(shape: GraphicFrame, data: KeyValueDataType):
    keys: KVKeys = data.get('keys', [])
    records = data.get('data', [])
    header_names: List[str] = [k.get('name', '') if isinstance(k, dict) else k for k in keys]
    data_keys: List[str] = [k.get('data_key') if isinstance(k, dict) else k for k in keys]
    for r_idx, row in enumerate(shape.table.rows):
        if r_idx == 0:
            for idx, name in enumerate(header_names):
                try:
                    replace_text_frame(row.cells[idx].text_frame, name)
                except Exception as e:
                    print(e)
        else:
            record = records[r_idx - 1]
            for idx, key in enumerate(data_keys):
                try:
                    replace_text_frame(row.cells[idx].text_frame, record.get(key))
                except Exception as e:
                    print(e)


def replace_table_data_raw(shape: GraphicFrame, data: KeyValueDataType):
    records = data.get('data', [])
    for r_idx, row in enumerate(shape.table.rows):
        try:
            record = records[r_idx]
            for c_idx, cell in enumerate(row.cells):
                replace_text_frame(cell.text_frame, record[c_idx])
        except Exception as e:
            print(e)


TABLE_DATA_TYPE_HANDLER = {
    "key_value": replace_table_data_key_value,
    "raw": replace_table_data_raw,
}


def replace_category_data(shape: GraphicFrame, data: CategoryDataType):
    chart_data = CategoryChartData()
    chart_data.categories = data.get('categories', [])
    for name, series_data in data.get('series', {}).items():
        chart_data.add_series(name, series_data)
    shape.chart.replace_data(chart_data)
    pass


CHART_DATA_TYPE_HANDLER: Dict[str, Callable[[GraphicFrame, ChartDataTypes], None]] = {
    "category_data": replace_category_data
}


def replace_chart_data(shape: GraphicFrame, chart_data: ChartDataType):
    data_type = chart_data['data_type']
    CHART_DATA_TYPE_HANDLER[data_type](shape, chart_data)
    if title := chart_data.get('title'):
        replace_text_frame(shape.chart.chart_title.text_frame, title)


def replace_data(slide, contents: dict, mode="replace"):
    for name, data in contents.items():
        shape = find_shape(slide, name) if mode == 'replace' else find_shape_by_slide_layout(slide, name)
        if shape:
            # shape.name = name # for
            if text_data := data.get('text'):
                replace_text_frame(shape.text_frame, text_data)
            elif table_data := data.get('table'):
                try:
                    data_type = table_data.get('data_type', 'raw')
                    TABLE_DATA_TYPE_HANDLER[data_type](shape, table_data)
                except Exception as e:
                    print(e)
            elif chart_data := data.get('chart'):
                try:
                    replace_chart_data(shape, chart_data)
                except Exception as e:
                    print(e)
